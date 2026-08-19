#!/usr/bin/env python3
"""Erzeugt die Abschlusspraesentation (Projekt I1) als PowerPoint (.pptx).

Bettet die Figuren aus ../paper/figures ein. Aufruf:
    ../.venv/bin/python build_praesentation.py
"""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Pfade ─────────────────────────────────────────────────────────────────────
HERE = Path(__file__).resolve().parent
FIG  = HERE.parent / "paper" / "figures"
OUT  = HERE / "Abschlusspraesentation_Plant_AE.pptx"
ROOT = HERE.parent.parent.parent  # Repo-Wurzel (projekt_i1)

PLANT_PHOTO = ROOT / "bilder" / "WhatsApp Video 2026-06-03 at 20.09.48_frame_00-00-27-16.png"
SCOPE_SPECTRUM = (HERE.parent / "data" / "hybrid_watering_experiment" / "20260622_224324"
                  / "snapshot_post_watering" / "plots" / "oscilloscope_spectrum.png")

# ── Farben ────────────────────────────────────────────────────────────────────
PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)   # dunkelblau
ACCENT  = RGBColor(0x2C, 0xA0, 0x2C)   # gruen (= CH3)
ALERT   = RGBColor(0xC0, 0x39, 0x2B)   # rot (Kernbefund)
INK     = RGBColor(0x22, 0x22, 0x22)
MUTED   = RGBColor(0x60, 0x66, 0x70)
LIGHT   = RGBColor(0xEE, 0xF2, 0xF6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

# ── Aufbau ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def textbox(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6):
    """lines: list of (text, size, color, bold, italic) or list thereof (paragraphs)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        segs = para if isinstance(para, list) else [para]
        for seg in segs:
            text, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            italic = seg[4] if len(seg) > 4 else False
            r = p.add_run(); r.text = text
            _set(r, size, color, bold, italic)
    return tb


def rect(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def bullets(slide, l, t, w, h, items, size=18, gap=10, color=INK):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        # it: (level, text, [color], [bold])
        lvl = it[0]; text = it[1]
        col = it[2] if len(it) > 2 and it[2] else color
        bold = it[3] if len(it) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap); p.space_before = Pt(0)
        bullet = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = bullet + text
        _set(r, size - lvl * 2, col, bold)
    return tb


def header(slide, kicker, title):
    """Kopfzeile fuer Inhaltsfolien."""
    rect(slide, 0, 0, SW, Inches(1.15), PRIMARY)
    rect(slide, 0, Inches(1.15), SW, Pt(3), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.14), Inches(12), Inches(0.35),
            [[(kicker.upper(), 12, RGBColor(0xBF, 0xD3, 0xE6), True)]])
    textbox(slide, Inches(0.55), Inches(0.42), Inches(12.2), Inches(0.68),
            [[(title, 27, WHITE, True)]])


def footer(slide, page):
    textbox(slide, Inches(0.55), Inches(7.06), Inches(9), Inches(0.35),
            [[("Plant AE Monitoring  ·  Systemcharakterisierung", 10, MUTED)]])
    textbox(slide, Inches(11.6), Inches(7.06), Inches(1.2), Inches(0.35),
            [[(f"{page}", 10, MUTED)]], align=PP_ALIGN.RIGHT)


def takeaway(slide, text, color=PRIMARY, top=Inches(6.15)):
    """Kernaussagen-Band unten."""
    bar = rect(slide, Inches(0.55), top, Inches(12.23), Inches(0.72), LIGHT)
    rect(slide, Inches(0.55), top, Pt(6), Inches(0.72), color)
    tb = slide.shapes.add_textbox(Inches(0.85), top, Inches(11.8), Inches(0.72))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Kernaussage:  "; _set(r, 15, color, True)
    r = p.add_run(); r.text = text; _set(r, 15, INK)
    return bar


def figure(slide, name, top=Inches(1.45), max_h=Inches(4.5), max_w=Inches(11.6)):
    path = FIG / name
    w, h = Image.open(path).size
    ar = w / h
    width = max_w
    height = Emu(int(width / ar))
    if height > max_h:
        height = max_h
        width = Emu(int(height * ar))
    left = Emu(int((SW - width) / 2))
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return top + height


def picture(slide, path, left, top, box_w, box_h):
    """Bild von beliebigem Pfad, in Box (box_w x box_h) eingepasst und zentriert."""
    w, h = Image.open(path).size
    ar = w / h
    width = box_w
    height = Emu(int(width / ar))
    if height > box_h:
        height = box_h
        width = Emu(int(height * ar))
    off_l = left + Emu(int((box_w - width) / 2))
    off_t = top + Emu(int((box_h - height) / 2))
    slide.shapes.add_picture(str(path), off_l, off_t, width=width, height=height)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# Folie 1 — Titel
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, PRIMARY)
rect(s, 0, Inches(4.55), SW, Pt(3.5), ACCENT)
textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.4),
        [[("ABSCHLUSSPRÄSENTATION  ·  PROJEKT I1  ·  SOMMERSEMESTER 2026",
           13, RGBColor(0xBF, 0xD3, 0xE6), True)]])
textbox(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(2.2),
        [[("Passive akustische Emissions-Überwachung", 40, WHITE, True)],
         [("von Pflanzenstress", 40, WHITE, True)],
         [("Systemcharakterisierung und experimentelle Validierung", 22,
           RGBColor(0xCF, 0xDE, 0xEA), False, True)]], space_after=8)
textbox(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.2),
        [[("Christoph Bellmann  ·  Julian Hübke", 20, WHITE, True)],
         [("Technische Hochschule Augsburg  ·  Systems Engineering", 15,
           RGBColor(0xCF, 0xDE, 0xEA))],
         [("christoph.bellmann@tha.de  ·  julianhuebke@tha.de", 13,
           RGBColor(0x9F, 0xB8, 0xCE))]],
        space_after=4)
notes(s, "Begruessung. Thema: Kann ein kostenguenstiges System akustische "
         "Emissionen von Pflanzen unter Wasserstress messen? Zunaechst mit Piezo-"
         "Sensoren, spaeter mit MEMS-Mikrofonen. Ich stelle Aufbau, "
         "Charakterisierung und die entscheidenden Experimente vor.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 2 — Motivation & Ziel
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Motivation & Ziel", "Warum akustische Emissionen von Pflanzen?")
bullets(s, Inches(0.6), Inches(1.5), Inches(6.15), Inches(4.5), [
    (0, "Pflanzen erzeugen akustische Emissionen (AE): winzige elastische "
        "Wellen aus Xylem-Kavitation und Zellwand-Mikrorissen.", INK),
    (0, "AE steigt unter Wasserstress → potenzieller Echtzeit-Proxy für "
        "die Bewässerungssteuerung.", INK),
    (0, "Klassische AE-Forschung nutzt teure resonante Ultraschallsensoren "
        ">100 kHz.", MUTED),
    (0, "Diese Arbeit: Lässt sich das mit kostengünstiger Hardware "
        "(Piezo + LM358 + Oszilloskop) messen?", PRIMARY, True),
], size=18, gap=14)
rect(s, Inches(7.05), Inches(1.55), Inches(5.75), Inches(4.35), LIGHT)
textbox(s, Inches(7.35), Inches(1.7), Inches(5.2), Inches(0.5),
        [[("Fünf Beiträge", 18, PRIMARY, True)]])
bullets(s, Inches(7.35), Inches(2.2), Inches(5.2), Inches(3.6), [
    (0, "Frequenz-Charakterisierung der 3-Kanal-Messkette", INK),
    (0, "Laufzeit-Analyse zur Sensorgeometrie", INK),
    (0, "Bewässerungsexperiment mit Drift-Korrektur", INK),
    (0, "Identifikation der limitierenden Rauschquellen", INK),
    (0, "Phase-3-Umbau auf MEMS-Mikrofone (Luft vs. Boden)", ACCENT, True),
], size=15, gap=11)
footer(s, 2)
notes(s, "AE = akustische Emission. Idee: Wenn Wasserstress die Kavitationsrate "
         "erhoeht, koennte man Trockenheit hoeren. Frage: reicht dafuer billige "
         "Hardware? Fuenf Beitraege gliedern den Vortrag – der fuenfte ist der "
         "Phase-3-Umbau auf MEMS-Mikrofone.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 3 — Messsystem (fig1)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Aufbau", "Das Messsystem: drei Piezo-Kanäle")
figure(s, "fig1_system.png", top=Inches(1.55), max_h=Inches(3.0), max_w=Inches(11.0))
bullets(s, Inches(0.7), Inches(4.75), Inches(12), Inches(1.3), [
    (0, "CH1 = Boden-Piezo im Wurzelbereich (Pflanzenkanal); CH3/CH4 = "
        "Referenz-Piezos an einem Stahlstab.  CH2 defekt.", INK),
    (0, "Verstärkung per LM358; Aufnahme mit Rigol DS1104Z: 500 kSa/s, "
        "300 000 Punkte pro Frame (0,6 s).", INK),
], size=16, gap=10)
takeaway(s, "Bewusst kostengünstige Kette – genau ihre Grenzen sollen "
            "quantifiziert werden.")
footer(s, 3)
notes(s, "Drei Kanaele: ein Sensor im Boden bei der Pflanze, zwei Referenzsensoren "
         "am Stahlstab fuer Laufzeit/Kalibrierung. LM358 ist ein Standard-Cent-"
         "Verstaerker. Aufnahme im Deep-Memory-Modus.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 4 — Realer Aufbau & Messsignal (Foto + Oszilloskop-Spektrum)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Aufbau", "Realer Messaufbau & Oszilloskop-Signal")
# linke Spalte: Foto der Pflanze
rect(s, Inches(0.55), Inches(1.55), Inches(5.9), Inches(0.42), LIGHT)
textbox(s, Inches(0.75), Inches(1.57), Inches(5.6), Inches(0.4),
        [[("Pflanze mit Boden-Piezo (CH1)", 15, PRIMARY, True)]])
picture(s, PLANT_PHOTO, Inches(0.55), Inches(2.1), Inches(5.9), Inches(3.9))
# rechte Spalte: Oszilloskop-Spektrum
rect(s, Inches(6.85), Inches(1.55), Inches(5.9), Inches(0.42), LIGHT)
textbox(s, Inches(7.05), Inches(1.57), Inches(5.6), Inches(0.4),
        [[("Gemessenes Spektrum (nach Bewässerung)", 15, ACCENT, True)]])
picture(s, SCOPE_SPECTRUM, Inches(6.85), Inches(2.1), Inches(5.9), Inches(3.9))
takeaway(s, "Realer Aufbau und ein tatsächlich aufgenommenes Spektrum – die "
            "Grundlage der folgenden Charakterisierung.")
footer(s, 4)
notes(s, "Links das reale Setup: die Pflanze mit dem Boden-Piezo (CH1). Rechts ein "
         "tatsaechlich aufgenommenes Oszilloskop-Spektrum aus dem "
         "Bewaesserungsexperiment. Damit wird der abstrakte Schaltplan der "
         "Vorfolie greifbar.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 5 — Methodik
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Vorgehen", "Messgrößen und Auswertung")
cols = [
    ("Aufnahme", [
        "Deep-Memory: 500 kSa/s, 0,6 s/Frame",
        "AC-Kopplung, 0,5 V/div",
        "Raspberry Pi Pico erzeugt bekannten 20–100 kHz-Chirp zur Kalibrierung",
    ], PRIMARY),
    ("Kennzahlen", [
        "Multi-Frame-Kohärenz γ² (Kanaltreue)",
        "MLE-Laufzeit τ̂ zwischen Sensoren",
        "Bandenergien in 5-kHz-Bins",
    ], ACCENT),
    ("Statistik", [
        "Wilcoxon-Test (gepaart, prä/post)",
        "FDR-Korrektur (Benjamini–Hochberg)",
        "Kendall τ (Trend) + Mann–Whitney (drift-korrigiert)",
    ], ALERT),
]
x = Inches(0.6)
for title, items, col in cols:
    rect(s, x, Inches(1.55), Inches(3.95), Inches(4.6), LIGHT)
    rect(s, x, Inches(1.55), Inches(3.95), Pt(6), col)
    textbox(s, x + Inches(0.25), Inches(1.75), Inches(3.5), Inches(0.5),
            [[(title, 19, col, True)]])
    bullets(s, x + Inches(0.2), Inches(2.4), Inches(3.6), Inches(3.5),
            [(0, it, INK) for it in items], size=15, gap=13)
    x += Inches(4.12)
footer(s, 5)
notes(s, "Kalibrierung mit bekanntem Chirp -> wir kennen das Eingangssignal und "
         "koennen die Kette vermessen. Kohaerenz misst, ob zwei Kanaele dasselbe "
         "sehen. Statistik ist bewusst streng, inkl. Drift-Kontrolle.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 6 — Bandbreite / Kohärenz (fig2)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Systemcharakterisierung (1/3)", "Nutzbares Frequenzband: 5–35 kHz")
figure(s, "fig2_bandwidth.png", top=Inches(1.55), max_h=Inches(4.2))
takeaway(s, "Kalibriertes Band 5–35 kHz (γ² ≥ 0,95); darüber bricht der "
            "LM358 ein. Das primäre Kavitationsband (100–500 kHz) ist unerreichbar.")
footer(s, 6)
notes(s, "Links: Kohaerenz nahe 1 bis 35 kHz -> beide Referenzsensoren sehen "
         "dasselbe, Kette ueberträgt sauber. Rechts: oberhalb 35 kHz faellt die "
         "Amplitude ab (LM358-Grenze). Peaks sind Stahlstab-Eigenmoden.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 7 — Laufzeit / Geometrie (fig3)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Systemcharakterisierung (2/3)", "Sensorgeometrie über Laufzeit")
figure(s, "fig3_tde.png", top=Inches(1.55), max_h=Inches(4.2))
takeaway(s, "Laufzeit ≈ 0 im Kalibrierband → Sensoren äquidistant. Der feste "
            "+1,95 µs-Versatz bei 9,17 kHz ist ein Resonanz-Phasenartefakt, "
            "keine Geometrie.")
footer(s, 7)
notes(s, "Links: im Nutzband liegt die Laufzeitdifferenz innerhalb ±0,3 µs -> "
         "Sensoren gleich weit entfernt. Rechts: an der Stahl-Resonanz ein "
         "konstanter 1,95-µs-Versatz, ueber alle Frames identisch -> Artefakt "
         "der Resonanz, nicht der Wegdifferenz.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 8 — Rauschgrenze (fig4)  — Kernfolie
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Systemcharakterisierung (3/3)", "Die Rauschgrenze — der entscheidende Punkt")
figure(s, "fig4_noise.png", top=Inches(1.5), max_h=Inches(4.2))
takeaway(s, "Der Bodenkanal (CH1) rauscht mit 137 mV RMS – dominiert durch "
            "1/f- und Netzeinkopplung. Jedes Pflanzensignal müsste diese Kurve "
            "überschreiten.", color=ALERT)
footer(s, 8)
notes(s, "Ohne anliegendes Signal aufgenommen: das ist reines Rauschen jedes "
         "Kanals. CH4 ist am leisesten (Verstaerker-Grundrauschen). CH1, der "
         "Pflanzenkanal, ist am lautesten. Diese Kurve ist die Nachweisschwelle.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 9 — Bewässerungsexperiment (fig5)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Experiment", "Bewässerung: gießassoziierter Verlauf, Ursache offen")
figure(s, "fig5_watering.png", top=Inches(1.5), max_h=Inches(4.2))
takeaway(s, "CH1 zeigt nach dem Gießen einen Pegelabfall mit langsamer Erholung. "
            "Ursache offen: Sensorankopplung, Substratänderung oder biologische "
            "Reaktion; kein eindeutiger AE-Nachweis.")
footer(s, 9)
notes(s, "Euphorbia leuconeura, einzelnes Bewaesserungsereignis. CH1 faellt beim "
         "Giessen um etwa 2,9 mV ab und erholt sich ueber rund 20 Minuten. Da die "
         "Verstaerker dauerhaft liefen und die Referenzkanaele keinen gemeinsamen "
         "Trend zeigen, ist thermische Verstaerkerdrift nicht belegt. Der Verlauf "
         "kann aus Sensorankopplung, Substrataenderung oder biologischer Reaktion "
         "stammen; ohne Rohwellenformen ist keine eindeutige AE-Zuordnung moeglich.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 10 — Positivkontrolle (fig7)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Kontrolle", "Positivkontrolle: die Pipeline funktioniert")
figure(s, "fig7_reference.png", top=Inches(1.5), max_h=Inches(4.2))
takeaway(s, "Ein bekannt eingespeistes Signal wird mit +52,6 dB Abstand klar "
            "erkannt (p < 10⁻¹²). Das Negativergebnis ist echt – kein "
            "Pipeline-Fehler.", color=ACCENT)
footer(s, 10)
notes(s, "Gegenprobe: speisen wir ein echtes Signal in CH1, trennt die Auswertung "
         "es eindeutig vom stillen Kanal. Damit ist bewiesen: findet die Pipeline "
         "nichts, ist wirklich nichts da – nicht die Auswertung ist blind.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 11 — Phase-3-Umbau auf MEMS & Salatpflanze (fig_mems_phase3)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Phase 3", "Umbau auf MEMS-Mikrofone & Salatpflanze")
bullets(s, Inches(0.6), Inches(1.45), Inches(12.2), Inches(1.6), [
    (0, "Aufbau umgezogen; CH3/CH4 von Piezo auf MEMS-Mikrofone umgerüstet "
        "(gleiche LM-Verstärkerstufen).", INK),
    (0, "Neue Pflanze seit 14.07.2026: Salat (Lactuca sativa); kein Home "
        "Assistant → keine automatische Bewässerung/Temperatur.", INK),
    (0, "MEMS-Kette direkt aus Rohdaten neu charakterisiert (Sweep 5–100 kHz, "
        "80 Frames; Rohdaten automatisch per git-LFS versioniert).", MUTED),
], size=15, gap=8)
figure(s, "fig_mems_phase3.png", top=Inches(3.15), max_h=Inches(2.7), max_w=Inches(12.2))
takeaway(s, "Kalibriertes MEMS-Band wieder 5–35 kHz. Boden-Rauschen dominiert von "
            "einem gemeinsamen ~5,36-kHz-Kamm: Netzteil-Tausch verschob ihn (−35 Hz), "
            "beseitigte ihn aber nicht → Quelle nicht isoliert.")
footer(s, 11)
notes(s, "Phase 3: neuer Standort, CH3/CH4 jetzt MEMS-Mikrofone statt Piezo, neue "
         "Pflanze (Salat). Die MEMS-Kette wurde mit demselben Verfahren wie die "
         "Piezo-Kette charakterisiert – nutzbares Band bleibt 5–35 kHz. Wichtig: der "
         "5,36-kHz-Stoerkamm ist NICHT eindeutig dem Netzteil zuzuordnen; ein A/B-"
         "Tausch verschob nur die Frequenz. Rohdaten werden automatisch per git-LFS "
         "committet – eine versionierte, nachvollziehbare Pipeline.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 12 — Phase-3 MEMS Befund (methods_overview)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Phase 3 – Befund", "MEMS in Luft & Boden: kein Pflanzen-AE")
picture(s, HERE.parent / "data" / "ch3_ch4_phase_shift_20260626" /
        "20260716_165231" / "methods_overview.png",
        Inches(0.55), Inches(1.5), Inches(6.5), Inches(4.4))
bullets(s, Inches(7.25), Inches(1.7), Inches(5.6), Inches(4.2), [
    (0, "Boden-Einführung hebt den Pegel (Median-RMS 80 → 303 mV), CH3/CH4 "
        "bleiben nahezu gleich.", INK),
    (0, "Dichter, kohärenter Linienkamm; stärkste Linie 32,59 kHz mit "
        "γ² = 0,997.", INK),
    (0, "GCC-PHAT = 0,000 µs → common-mode, nicht lokalisierbar (keine "
        "Geometrie-Information).", ALERT, True),
    (0, "CH3-Sättigung in Frame 13; kein Band übersteht die FDR-Korrektur "
        "(0/32).", INK),
], size=15, gap=12)
takeaway(s, "Beide MEMS funktionieren, aber weder Luft noch Boden liefern Pflanzen-"
            "AE – der Boden verstärkt vor allem ein kohärentes Störfeld.", color=ALERT)
footer(s, 12)
notes(s, "Zwei passive MEMS-Kampagnen: Luft-Referenz und beide Mikrofone im Boden. "
         "Der Boden hebt den Pegel stark, aber die Kanaele sind fast identisch und "
         "die Laufzeit ist exakt null – das ist eine gemeinsame Einkopplung, kein "
         "ortbares Pflanzensignal. Dazu ein Saettigungsfehler auf CH3. Kein "
         "statistisch belastbares Pflanzen-AE.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 13 — Fazit
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Fazit", "Was wir gelernt haben")
bullets(s, Inches(0.7), Inches(1.5), Inches(12.1), Inches(4.2), [
    (0, "Kalibriertes Messband 5–35 kHz (Piezo und MEMS); Sensorgeometrie "
        "über Laufzeit bestätigt.", INK),
    (0, "Kein Pflanzen-AE im Bewässerungsexperiment – nach strenger "
        "Drift-Korrektur (p = 0,86).", INK),
    (0, "Phase-3-MEMS (Luft & Boden): ebenfalls kein Pflanzen-AE; Boden zeigt "
        "nur ein kohärentes Störfeld (32,59 kHz) + CH3-Sättigung.", ACCENT, True),
    (0, "Positivkontrolle belegt: die Auswertung ist funktionsfähig, das "
        "Negativergebnis ist belastbar.", INK),
    (0, "Hauptlimit: LM358-Rauschgrenze und common-mode-Einkopplung – 20–40 dB "
        "über erwarteten AE-Amplituden.", ALERT, True),
    (0, "Zusätzlich schließt die Bandbreite (<35 kHz) das Kavitationsband "
        "(100–500 kHz) aus.", INK),
], size=17, gap=11)
footer(s, 13)
notes(s, "Ehrliches Ergebnis über beide Sensorgenerationen: weder Piezo noch MEMS "
         "weisen Pflanzen-AE nach. Aber wir wissen jetzt genau WARUM – "
         "Rauschgrenze, common-mode-Einkopplung, Bandbreite – und das ist verwertbar.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 14 — Ausblick
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Ausblick", "Weg zur Nachweisbarkeit")
rows = [
    ("1", "Rauscharmer Vorverstärker (e_n ≈ 4 nV/√Hz, z. B. OPA2134)", "≈ 18 dB"),
    ("2", "Hardware-Hochpass 1 kHz am Verstärkereingang (1/f + Netz)", "≈ 20 dB"),
    ("3", "Resonanter 150-kHz-Piezo + 5 MSa/s-ADC → Kavitationsband", "neuer Bereich"),
    ("4", "Mechanische Isolation (schwingungsentkoppelter Sockel)", "Störungen ↓"),
    ("5", "MEMS: Geometrie dokumentieren, common-mode entkoppeln, Impuls-"
          "Positivkontrolle", "Phase 3"),
]
y = Inches(1.55)
for num, txt, gain in rows:
    rect(s, Inches(0.7), y, Inches(0.6), Inches(0.72), PRIMARY)
    textbox(s, Inches(0.7), y, Inches(0.6), Inches(0.72),
            [[(num, 20, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.4), y, Inches(9.1), Inches(0.72), LIGHT)
    textbox(s, Inches(1.65), y, Inches(8.7), Inches(0.72),
            [[(txt, 15, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(10.6), y, Inches(2.2), Inches(0.72), RGBColor(0xE2, 0xEC, 0xF3))
    textbox(s, Inches(10.6), y, Inches(2.2), Inches(0.72),
            [[(gain, 14, PRIMARY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.86)
takeaway(s, "Schritte 1 + 2 bringen CH1 unter 5 mV RMS; Schritt 5 macht künftige "
            "MEMS-Läufe erst interpretierbar.")
footer(s, 14)
notes(s, "Konkrete, priorisierte Massnahmen mit geschaetztem Gewinn. Schritte 1 und "
         "2 sind billig und wirksam. Schritt 5 adressiert die MEMS-spezifischen "
         "Luecken: ohne dokumentierte Sensorabstaende, common-mode-Entkopplung und "
         "eine zeitlich verankerte Positivkontrolle bleiben passive MEMS-Kandidaten "
         "nicht der Pflanze zuzuordnen.")

# ══════════════════════════════════════════════════════════════════════════════
# Folie 15 — Abschluss
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, PRIMARY)
rect(s, 0, Inches(2.35), SW, Pt(3.5), ACCENT)
textbox(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.0),
        [[("Vielen Dank für Ihre Aufmerksamkeit", 34, WHITE, True)]])
bullets(s, Inches(0.95), Inches(3.0), Inches(11.4), Inches(3.0), [
    (0, "Zwei Sensorgenerationen charakterisiert (Piezo & MEMS): nutzbares "
        "Band 5–35 kHz.", RGBColor(0xDD, 0xE8, 0xF1)),
    (0, "Kein Pflanzen-AE nachweisbar – begründet durch Rauschgrenze und "
        "common-mode-Einkopplung, nicht durch die Auswertung.",
        RGBColor(0xDD, 0xE8, 0xF1)),
    (0, "Klarer, priorisierter Weg zur Nachweisbarkeit aufgezeigt.",
        RGBColor(0xDD, 0xE8, 0xF1)),
], size=18, gap=14)
textbox(s, Inches(0.95), Inches(6.15), Inches(11), Inches(0.9),
        [[("Christoph Bellmann · Julian Hübke  ·  THA Augsburg", 15,
           RGBColor(0xBF, 0xD3, 0xE6))],
         [("Daten, Code und Paper: src/characterization/  (Rohdaten via git-LFS)", 13,
           RGBColor(0x9F, 0xB8, 0xCE))]], space_after=4)
notes(s, "Zusammenfassung in drei Saetzen ueber beide Sensorgenerationen. "
         "Offen fuer Fragen.")

# ── Speichern ─────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"gespeichert: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} Folien)")
